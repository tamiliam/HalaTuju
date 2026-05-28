"""
Generate the student-onboarding PowerPoint from the canonical Markdown content.

Run from anywhere:  python docs/scholarship/_build_deck.py
Output:             docs/scholarship/onboarding-deck.pptx

Why this exists:
- The Markdown deck (onboarding-deck.md) is the source of truth and the easy
  place to edit copy.
- The PowerPoint output is what the user presents (and can re-edit in PPT).
- python-pptx builds the deck programmatically — no Stitch / paid call needed.
- Visual identity matches HalaTuju: Lexend (Calibri fallback), brand blue
  #137fec, white background, generous spacing.

Delete this script after a successful build if you don't plan to regenerate;
otherwise keep it next to the .md for one-line re-builds.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
STITCH = HERE / 'stitch'
OUT = HERE / 'onboarding-deck.pptx'

# ── Brand ─────────────────────────────────────────────────────────────────
BRAND_BLUE = RGBColor(0x13, 0x7F, 0xEC)
TEXT_DARK = RGBColor(0x11, 0x1c, 0x2d)
TEXT_MUTED = RGBColor(0x5c, 0x5f, 0x61)
BG_SOFT = RGBColor(0xF0, 0xF3, 0xFF)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_TEXT = RGBColor(0x92, 0x40, 0x0E)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_TEXT = RGBColor(0x16, 0x65, 0x34)
RULE = RGBColor(0xE2, 0xE8, 0xF0)

FONT_HEAD = 'Lexend'   # falls back to Calibri/Segoe UI gracefully
FONT_BODY = 'Lexend'

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def _para_runs(p, runs, size=14, bold=False, color=TEXT_DARK, font=FONT_BODY,
               align=PP_ALIGN.LEFT):
    """Set a paragraph's text from a list of (text, **fmt_overrides) tuples."""
    p.alignment = align
    p.text = ''  # reset
    first = True
    for text, *fmt in runs:
        overrides = fmt[0] if fmt else {}
        if first:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
            first = False
        else:
            run = p.add_run()
            run.text = text
        f = run.font
        f.name = overrides.get('font', font)
        f.size = Pt(overrides.get('size', size))
        f.bold = overrides.get('bold', bold)
        f.color.rgb = overrides.get('color', color)


def _set_para(p, text, *, size=14, bold=False, color=TEXT_DARK, font=FONT_BODY,
              align=PP_ALIGN.LEFT, space_after=4):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    tb.text_frame.margin_left = Pt(0)
    tb.text_frame.margin_right = Pt(0)
    tb.text_frame.margin_top = Pt(0)
    tb.text_frame.margin_bottom = Pt(0)
    return tb.text_frame


def _add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def _add_page_number(slide, n, total):
    tf = _add_textbox(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3))
    p = tf.paragraphs[0]
    _set_para(p, f'{n} / {total}', size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def _slide_title(slide, eyebrow, title):
    """Eyebrow tag + slide title block at the top."""
    if eyebrow:
        tf = _add_textbox(slide, MARGIN, Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.35))
        p = tf.paragraphs[0]
        _set_para(p, eyebrow, size=12, bold=True, color=BRAND_BLUE, font=FONT_HEAD,
                  align=PP_ALIGN.LEFT)
    tf = _add_textbox(slide, MARGIN, Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.9))
    p = tf.paragraphs[0]
    _set_para(p, title, size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)


def _h_rule(slide, y):
    from pptx.enum.shapes import MSO_SHAPE
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, y, SLIDE_W - 2 * MARGIN, Pt(0.75),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()


def new_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


# ── Build ─────────────────────────────────────────────────────────────────


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 7

    # ── Slide 1: Title + welcome + HalaTuju intro ─────────────────────────
    s = new_slide(prs)
    # Big title
    tf = _add_textbox(s, MARGIN, Inches(0.9), SLIDE_W - 2 * MARGIN, Inches(1.6))
    p = tf.paragraphs[0]
    _set_para(p, 'B40 Education Assistance', size=44, bold=True, color=TEXT_DARK,
              font=FONT_HEAD, space_after=6)
    p2 = tf.add_paragraph()
    _set_para(p2, 'A community pilot — continuing your studies at Malaysia’s public colleges and universities',
              size=20, color=BRAND_BLUE, font=FONT_HEAD, space_after=4)

    # Soft intro box
    box = _add_rect(s, MARGIN, Inches(3.0), SLIDE_W - 2 * MARGIN, Inches(3.6), BG_SOFT)
    tf = _add_textbox(s, MARGIN + Inches(0.35), Inches(3.2), SLIDE_W - 2 * MARGIN - Inches(0.7), Inches(3.2))
    para = tf.paragraphs[0]
    _set_para(para, 'About this programme', size=14, bold=True, color=BRAND_BLUE, font=FONT_HEAD, space_after=2)
    p = tf.add_paragraph()
    _set_para(p,
              'A small, community-supported effort helping students from low-income households continue at Malaysia’s public institutions — backed by fellow Malaysians who want to see them succeed. There is no big organisation behind this yet; it is a pilot, run by people who care.',
              size=14, color=TEXT_DARK, space_after=10)
    p = tf.add_paragraph()
    _set_para(p, 'About HalaTuju', size=14, bold=True, color=BRAND_BLUE, font=FONT_HEAD, space_after=2)
    p = tf.add_paragraph()
    _set_para(p,
              'HalaTuju is a Malaysian education-pathway platform that helps SPM/STPM students find courses they qualify for. This assistance programme runs on HalaTuju — you will create or log in to your HalaTuju account to apply.',
              size=14, color=TEXT_DARK, space_after=10)
    p = tf.add_paragraph()
    _set_para(p, 'About this session', size=14, bold=True, color=BRAND_BLUE, font=FONT_HEAD, space_after=2)
    p = tf.add_paragraph()
    _set_para(p,
              'You already raised your hand on our Google form (thank you). The Google form told us who you are; today we will walk through the formal application that gives us the verifiable detail and consent we need before we can introduce you to a sponsor. Everything you share stays confidential — we never pass it on without your permission.',
              size=14, color=TEXT_DARK)
    _add_page_number(s, 1, TOTAL)

    # ── Slide 2: 8-step journey ───────────────────────────────────────────
    s = new_slide(prs)
    _slide_title(s, 'The /scholarship page', 'Your 8-step journey')
    tf = _add_textbox(s, MARGIN, Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(0.3))
    _set_para(tf.paragraphs[0], 'Rolling basis — no fixed deadline.', size=12, color=TEXT_MUTED)

    steps = [
        ('1', 'Apply',                'Sign in with Google, confirm your HalaTuju profile, submit.',                     'About 5–10 minutes'),
        ('2', 'We confirm',           'An email acknowledging your application.',                                         'Same day'),
        ('3', 'Shortlisting',         'We check your results and household income against the criteria and let you know.','Within 48 hours'),
        ('4', 'Complete your profile','If shortlisted, share a few more details + upload documents (IC, results slip, proof of income).', 'A few minutes when ready'),
        ('5', 'A short interview',    'Brief phone call (~20 min) in your preferred language. To get to know you, not test you.', 'After your profile is complete'),
        ('6', 'Ready for sponsors',   'With your written consent, we prepare a confidential profile.',                    'After the interview'),
        ('7', 'Matched and awarded',  'Once a sponsor chooses to support you, we arrange the assistance and administer the funds together.', 'Up to ~2 months from application'),
        ('8', 'Staying supported',    'Each semester, upload your latest results + a brief progress note. Support continues as long as you are progressing.', 'Throughout your studies'),
    ]
    rows = len(steps)
    top = Inches(2.1)
    row_h = (SLIDE_H - top - Inches(0.7)) / rows
    col_n = Inches(0.5)
    col_title = Inches(2.3)
    col_desc = Inches(7.8)
    col_when = Inches(2.4)

    for i, (n, title, desc, when) in enumerate(steps):
        y = top + i * row_h
        # number badge
        from pptx.enum.shapes import MSO_SHAPE
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y + Pt(2), Inches(0.32), Inches(0.32))
        badge.fill.solid(); badge.fill.fore_color.rgb = BRAND_BLUE
        badge.line.fill.background()
        tfb = badge.text_frame
        tfb.margin_left = Pt(0); tfb.margin_right = Pt(0); tfb.margin_top = Pt(0); tfb.margin_bottom = Pt(0)
        _set_para(tfb.paragraphs[0], n, size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                  font=FONT_HEAD, align=PP_ALIGN.CENTER)
        # title
        tf = _add_textbox(s, MARGIN + Inches(0.45), y, col_title, row_h)
        _set_para(tf.paragraphs[0], title, size=13, bold=True, color=TEXT_DARK, font=FONT_HEAD)
        # desc
        tf = _add_textbox(s, MARGIN + Inches(0.45) + col_title, y, col_desc, row_h)
        _set_para(tf.paragraphs[0], desc, size=11, color=TEXT_DARK)
        # when
        tf = _add_textbox(s, MARGIN + Inches(0.45) + col_title + col_desc, y, col_when, row_h)
        _set_para(tf.paragraphs[0], when, size=11, bold=True, color=BRAND_BLUE, font=FONT_HEAD, align=PP_ALIGN.RIGHT)

    _add_page_number(s, 2, TOTAL)

    # ── Slide 3: Eligibility ──────────────────────────────────────────────
    s = new_slide(prs)
    _slide_title(s, '/scholarship → Can I apply?', 'You are eligible if all of these are true')

    eligibility = [
        ('Malaysian citizen',          'who has just completed, or is completing, SPM or STPM.'),
        ('Low-income household',       'combined monthly income at or below RM5,860 (the national B40 threshold, DOSM 2024).'),
        ('Solid academic record',      'at least 5 A’s in SPM, or a PNGK of 3.0 or above in STPM.'),
        ('Continuing at a Malaysian public institution', 'Matrikulasi, Foundation/Asasi, STPM, IPTA, polytechnic or ILKA.'),
        ('Reachable for a short interview', 'usually a 20-minute phone call in English, Bahasa Malaysia, or Tamil.'),
        ('Willing to share progress',  'semester results + a brief update so support can continue.'),
    ]
    top = Inches(2.0)
    row_h = Inches(0.55)
    for i, (head, tail) in enumerate(eligibility):
        y = top + i * row_h
        # bullet number
        tf = _add_textbox(s, MARGIN, y, Inches(0.35), row_h)
        _set_para(tf.paragraphs[0], f'{i+1}.', size=14, bold=True, color=BRAND_BLUE, font=FONT_HEAD)
        # head + tail
        tf = _add_textbox(s, MARGIN + Inches(0.35), y, SLIDE_W - 2 * MARGIN - Inches(0.35), row_h)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run1 = p.runs[0] if p.runs else p.add_run()
        run1.text = head + ' — '
        run1.font.name = FONT_BODY; run1.font.size = Pt(14); run1.font.bold = True
        run1.font.color.rgb = TEXT_DARK
        run2 = p.add_run()
        run2.text = tail
        run2.font.name = FONT_BODY; run2.font.size = Pt(14); run2.font.color.rgb = TEXT_DARK

    # callout
    cy = top + len(eligibility) * row_h + Inches(0.15)
    _add_rect(s, MARGIN, cy, SLIDE_W - 2 * MARGIN, Inches(1.05), BG_SOFT)
    tf = _add_textbox(s, MARGIN + Inches(0.25), cy + Inches(0.12), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.85))
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.text = 'Close, but not exactly? '
    run.font.name = FONT_HEAD; run.font.size = Pt(14); run.font.bold = True
    run.font.color.rgb = BRAND_BLUE
    run2 = p.add_run()
    run2.text = 'Apply anyway. Our bar is a little more generous than the headline — especially on grades. Applying lets us assess your eligibility; it does not guarantee assistance.'
    run2.font.name = FONT_BODY; run2.font.size = Pt(13); run2.font.color.rgb = TEXT_DARK

    _add_page_number(s, 3, TOTAL)

    # ── Slide 4: /apply + truthfulness ────────────────────────────────────
    s = new_slide(prs)
    _slide_title(s, 'Step 1 — halatuju.xyz/scholarship/apply', 'The formal application form')
    tf = _add_textbox(s, MARGIN, Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(0.3))
    _set_para(tf.paragraphs[0], '5 short sections, save as you go.', size=12, color=TEXT_MUTED)

    sections = [
        ('About Me', 'Name, NRIC (one-time, pre-filled at sign-in), contact, school, preferred call language'),
        ('My Family', 'Household income, household size, who supports the family, STR/JKM if applicable'),
        ('My Results', 'Your SPM / STPM grades (re-uses what you have already entered on HalaTuju)'),
        ('My Plans', 'What you intend to study, where, and how confident you are (it is fine to be still deciding)'),
        ('Support', 'Other scholarships you are pursuing, anything else you would like us to know, consent to contact'),
    ]
    top = Inches(2.05)
    row_h = Inches(0.45)
    for i, (head, tail) in enumerate(sections):
        y = top + i * row_h
        tf = _add_textbox(s, MARGIN, y, Inches(2.4), row_h)
        _set_para(tf.paragraphs[0], head, size=13, bold=True, color=TEXT_DARK, font=FONT_HEAD)
        tf = _add_textbox(s, MARGIN + Inches(2.4), y, SLIDE_W - 2 * MARGIN - Inches(2.4), row_h)
        _set_para(tf.paragraphs[0], tail, size=12, color=TEXT_DARK)

    # Truthfulness callout (amber)
    cy = top + len(sections) * row_h + Inches(0.15)
    _add_rect(s, MARGIN, cy, SLIDE_W - 2 * MARGIN, Inches(1.5), AMBER_BG)
    tf = _add_textbox(s, MARGIN + Inches(0.25), cy + Inches(0.12), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(1.3))
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.text = '⚠  Truthfulness matters.  '
    run.font.name = FONT_HEAD; run.font.size = Pt(14); run.font.bold = True
    run.font.color.rgb = AMBER_TEXT
    run2 = p.add_run()
    run2.text = 'We later verify what you say against your documents (MyKad, results slip, income proof) and at the interview. '
    run2.font.name = FONT_BODY; run2.font.size = Pt(13); run2.font.color.rgb = TEXT_DARK
    run3 = p.add_run()
    run3.text = 'Mismatches will likely disqualify your application. '
    run3.font.name = FONT_BODY; run3.font.size = Pt(13); run3.font.bold = True; run3.font.color.rgb = AMBER_TEXT
    run4 = p.add_run()
    run4.text = 'If you are unsure of an exact number, be honest about your uncertainty in the open notes rather than guessing. We can work with “I’m not sure”; we cannot work with information that does not match.'
    run4.font.name = FONT_BODY; run4.font.size = Pt(13); run4.font.color.rgb = TEXT_DARK

    # After-submit footer
    fy = cy + Inches(1.6)
    tf = _add_textbox(s, MARGIN, fy, SLIDE_W - 2 * MARGIN, Inches(0.6))
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.text = 'After submit (Steps 2 → 3):  '
    run.font.name = FONT_HEAD; run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = BRAND_BLUE
    run2 = p.add_run()
    run2.text = 'a same-day acknowledgement email arrives; within 48 hours you will hear whether you have been shortlisted (link to the next page) or not this round (a warm note explaining why). Please check your spam folder.'
    run2.font.name = FONT_BODY; run2.font.size = Pt(12); run2.font.color.rgb = TEXT_DARK

    _add_page_number(s, 4, TOTAL)

    # ── Slide 5: /application + tell us as much as possible ───────────────
    s = new_slide(prs)
    _slide_title(s, 'Step 4 — halatuju.xyz/scholarship/application', 'Five small tabs, no fixed time limit')

    tabs = [
        ('1. Quiz', 'The HalaTuju course-fit quiz (re-uses your previous answers).'),
        ('2. Your story', 'About your family + about you. Mostly optional. Write in BM, English, or Tamil.'),
        ('3. Funding', 'Tick the categories support would help with (living, transport, accommodation, books, device). Plus a short note. Assistance is up to RM3,000 — could be less, depending on funds and need.'),
        ('4. Documents', 'Required: IC + results slip. Optional: one income proof (STR / salary slip / EPF), utility bills, statement of intent, offer letter, photo. The IC is auto-checked against your details — your photo isn’t kept at Google.'),
        ('5. Consent', 'Your formal permission to share your confidential profile with potential sponsors. You can withdraw it at any time.'),
    ]
    top = Inches(1.85)
    row_h = Inches(0.50)
    for i, (head, tail) in enumerate(tabs):
        y = top + i * row_h
        tf = _add_textbox(s, MARGIN, y, Inches(1.7), row_h)
        _set_para(tf.paragraphs[0], head, size=12, bold=True, color=BRAND_BLUE, font=FONT_HEAD)
        tf = _add_textbox(s, MARGIN + Inches(1.7), y, Inches(6.5), row_h)
        _set_para(tf.paragraphs[0], tail, size=11, color=TEXT_DARK)

    # Embed Documents + Whats-next screenshots on the right
    docs_img = STITCH / 's4-documents-built.png'
    next_img = STITCH / 's5-what-happens-next.png'
    img_top = Inches(1.85)
    img_w = Inches(1.8)
    if docs_img.exists():
        s.shapes.add_picture(str(docs_img), Inches(9.05), img_top, width=img_w)
    if next_img.exists():
        s.shapes.add_picture(str(next_img), Inches(11.05), img_top, width=img_w)
    cap = _add_textbox(s, Inches(9.05), img_top + Inches(4.8), Inches(3.8), Inches(0.3))
    _set_para(cap.paragraphs[0], 'Documents tab    →    "You’re all set!" + What happens next',
              size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Green "Tell us more" callout
    cy = top + len(tabs) * row_h + Inches(0.20)
    _add_rect(s, MARGIN, cy, Inches(8.3), Inches(1.05), GREEN_BG)
    tf = _add_textbox(s, MARGIN + Inches(0.25), cy + Inches(0.12), Inches(7.9), Inches(0.85))
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.text = '\U0001F4A1  Tell us as much as you can.  '
    run.font.name = FONT_HEAD; run.font.size = Pt(13); run.font.bold = True
    run.font.color.rgb = GREEN_TEXT
    run2 = p.add_run()
    run2.text = '"Your story" and the funding note are your chance to put a person behind the numbers — the more we know, the better we can advocate for you to sponsors. Write in your own words, in your most comfortable language. Same truthfulness rule applies: '
    run2.font.name = FONT_BODY; run2.font.size = Pt(12); run2.font.color.rgb = TEXT_DARK
    run3 = p.add_run()
    run3.text = 'honest uncertainty is fine; mismatches will likely disqualify.'
    run3.font.name = FONT_BODY; run3.font.size = Pt(12); run3.font.bold = True; run3.font.color.rgb = AMBER_TEXT

    _add_page_number(s, 5, TOTAL)

    # ── Slide 6: Steps 5–8 + Important things ────────────────────────
    s = new_slide(prs)
    _slide_title(s, 'After your profile is complete', 'Steps 5 to 8 — and the things to keep in mind')

    # Left column: Steps 5-8
    col_w = Inches(6.0)
    top = Inches(1.85)
    step_blocks = [
        ('Step 5  —  The phone call',     '~20 minutes, your preferred language. A warm conversation. We may also cross-check a few things — please be straightforward about anything you weren’t sure of in writing.'),
        ('Step 6  —  Profile for sponsors', 'A short, confidential profile based only on what you told us. A coordinator reviews it before any sponsor sees it.'),
        ('Step 7  —  Matched and awarded',  'Once a sponsor chooses to support you, we arrange the assistance and administer the funds together. From application to award typically takes up to two months.'),
        ('Step 8  —  Staying supported',    'Each semester you upload your latest results + a short progress note. As long as you are progressing, support continues through your studies.'),
    ]
    by = top
    for head, body in step_blocks:
        tf = _add_textbox(s, MARGIN, by, col_w, Inches(0.32))
        _set_para(tf.paragraphs[0], head, size=14, bold=True, color=BRAND_BLUE, font=FONT_HEAD, space_after=2)
        tf = _add_textbox(s, MARGIN, by + Inches(0.32), col_w, Inches(0.9))
        _set_para(tf.paragraphs[0], body, size=11, color=TEXT_DARK)
        by += Inches(1.20)

    # Right column: Important things plainly
    right_x = Inches(7.0)
    tf = _add_textbox(s, right_x, top, Inches(5.7), Inches(0.32))
    _set_para(tf.paragraphs[0], 'A few important things, plainly', size=14, bold=True,
              color=TEXT_DARK, font=FONT_HEAD, space_after=4)

    important = [
        ('\U0001F193  No fee, ever.', 'Free to apply, free to be in.'),
        ('\U0001F512  Confidential.', 'Your information goes only to our small coordinator team and (with your consent) to potential sponsors.'),
        ('✋  Soft signals, not hard blocks.', 'The IC auto-check is a hint — a human always reviews your documents before any decision.'),
        ('⚠  Truthfulness disqualifies.', 'Mismatches between what you say and what your documents show will likely end your application. Be honest about uncertainty.'),
        ('\U0001F468‍\U0001F469‍\U0001F467  Under 18?', 'A parent/guardian co-signs the consent. We will guide you through it.'),
        ('\U0001F1F2\U0001F1FE  Public institutions only,', 'this round — Matrikulasi, Foundation/Asasi, STPM, IPTA, polytechnic, ILKA.'),
    ]
    iy = top + Inches(0.5)
    for head, tail in important:
        tf = _add_textbox(s, right_x, iy, Inches(5.7), Inches(0.7))
        p = tf.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = head + ' '
        run.font.name = FONT_HEAD; run.font.size = Pt(12); run.font.bold = True
        run.font.color.rgb = TEXT_DARK
        run2 = p.add_run()
        run2.text = tail
        run2.font.name = FONT_BODY; run2.font.size = Pt(11); run2.font.color.rgb = TEXT_DARK
        iy += Inches(0.72)

    _add_page_number(s, 6, TOTAL)

    # ── Slide 7: Questions ────────────────────────────────────────────────
    s = new_slide(prs)
    tf = _add_textbox(s, MARGIN, Inches(1.6), SLIDE_W - 2 * MARGIN, Inches(1.4))
    _set_para(tf.paragraphs[0], 'Questions?', size=54, bold=True, color=TEXT_DARK,
              font=FONT_HEAD, align=PP_ALIGN.CENTER, space_after=6)
    p = tf.add_paragraph()
    _set_para(p, 'We’re glad you’re here.', size=20, color=BRAND_BLUE,
              font=FONT_HEAD, align=PP_ALIGN.CENTER)

    qs = [
        ('\U0001F4E8',  'Email us at the address on halatuju.xyz/scholarship — we reply within 2 working days.'),
        ('\U0001F4DE',  'The interview is a great time to ask anything you weren’t sure about.'),
        ('\U0001F91D',  'We’re a small group of Malaysians helping fellow Malaysians. Not lenders — a community pilot. Welcome.'),
    ]
    top = Inches(4.0)
    for i, (icon, text) in enumerate(qs):
        y = top + i * Inches(0.7)
        tf = _add_textbox(s, Inches(2.5), y, Inches(0.6), Inches(0.55))
        _set_para(tf.paragraphs[0], icon, size=18, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        tf = _add_textbox(s, Inches(3.1), y + Inches(0.05), Inches(8.5), Inches(0.55))
        _set_para(tf.paragraphs[0], text, size=14, color=TEXT_DARK)

    _add_page_number(s, 7, TOTAL)

    prs.save(OUT)
    print(f'OK — wrote {OUT}  ({OUT.stat().st_size // 1024} KB, 7 slides)')


if __name__ == '__main__':
    build()
